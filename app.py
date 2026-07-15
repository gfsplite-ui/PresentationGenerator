import json
import re
import os
import time
import io
import ast
import unicodedata
import waitress
from datetime import date
import tempfile
import httpx
import requests
from pptx import Presentation
from pptx.util import Pt, Inches, Emu, Cm
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from gigachat import GigaChat
from PIL import Image, ImageDraw, ImageFont
from flask import Flask, render_template_string, request, send_file, jsonify

import urllib3

urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# Импорт конфигурации (если config.py рядом)
try:
    from config import (
        GIGACHAT_CREDENTIALS,
        TEMPLATE_PATH,
        TITLE_LAYOUT_INDEX,
        CONTENT_LAYOUT_INDEX,
        ADD_IMAGES,
        SOFT_EDGE_RADIUS,
    )
except ImportError:
    # Для теста можно задать вручную, но лучше создать config.py
    GIGACHAT_CREDENTIALS = "ваш_ключ"
    TEMPLATE_PATH = "company_template.pptx"
    TITLE_LAYOUT_INDEX = 0
    CONTENT_LAYOUT_INDEX = 13
    ADD_IMAGES = True
    SOFT_EDGE_RADIUS = 10

giga = GigaChat(
    credentials=GIGACHAT_CREDENTIALS, scope="GIGACHAT_API_PERS", verify_ssl_certs=False
)


# ─── Вспомогательные функции ─────────────────────────────────────
def safe_filename(text):
    forbidden_chars = r'<>:"/\|?*'
    cleaned = "".join(c for c in text if c not in forbidden_chars)
    cleaned = cleaned.strip().rstrip(".")
    return cleaned if cleaned else "presentation"


def clean_json_string(s):
    allowed_controls = {10, 13, 9}
    cleaned = []
    for ch in s:
        cp = ord(ch)
        if cp < 0x20:
            if cp in allowed_controls:
                cleaned.append(ch)
            else:
                cleaned.append(" ")
        elif unicodedata.category(ch).startswith("C"):
            continue
        else:
            cleaned.append(ch)
    return "".join(cleaned).strip()


def extract_json_from_text(text):
    start = text.find("{")
    if start == -1:
        return text
    end = text.rfind("}")
    if end == -1:
        return text[start:]
    return text[start : end + 1]


def repair_json_structure(json_str):
    braces = 0
    brackets = 0
    in_string = False
    escape = False
    for ch in json_str:
        if escape:
            escape = False
            continue
        if ch == "\\":
            escape = True
            continue
        if ch == '"':
            in_string = not in_string
        if in_string:
            continue
        if ch == "{":
            braces += 1
        elif ch == "}":
            braces -= 1
        elif ch == "[":
            brackets += 1
        elif ch == "]":
            brackets -= 1
    repaired = json_str.rstrip()
    repaired += "]" * max(0, brackets)
    repaired += "}" * max(0, braces)
    return repaired


def parse_json_safe(text):
    json_str = extract_json_from_text(text)
    if not json_str:
        return None
    try:
        return json.loads(json_str)
    except json.JSONDecodeError:
        pass
    repaired = repair_json_structure(json_str)
    try:
        return json.loads(repaired)
    except json.JSONDecodeError:
        pass
    if '"' not in json_str:
        try:
            replaced = json_str.replace("'", '"')
            return json.loads(replaced)
        except (json.JSONDecodeError, Exception):
            pass
    try:
        return ast.literal_eval(json_str)
    except (ValueError, SyntaxError):
        pass
    return None


def normalize_title(title):
    if not title:
        return title
    words = title.split()
    if not words:
        return title
    new_words = [words[0].capitalize()]
    for word in words[1:]:
        if word.isupper() and len(word) > 1:
            new_words.append(word)
        else:
            new_words.append(word.lower())
    return " ".join(new_words)


def clean_subtitle(text):
    if not text:
        return ""
    text = re.sub(r"^[#]+\s*", "", text)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"\n+", " ", text)
    text = text.strip()
    if len(text) > 120:
        text = text[:120].rsplit(" ", 1)[0] + "…"
    return text


def delete_all_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def add_soft_edge(picture, radius_pt=10):
    radius_emu = int(radius_pt * 12700)
    soft_edge_xml = (
        f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:softEdge rad="{radius_emu}"/>'
        f"</a:effectLst>"
    )
    pic_xml = picture._element
    sp_pr = pic_xml.find(qn("p:spPr"))
    if sp_pr is None:
        sp_pr = pic_xml.makeelement(qn("p:spPr"), {})
        pic_xml.insert(0, sp_pr)
    existing = sp_pr.find(qn("a:effectLst"))
    if existing is not None:
        sp_pr.remove(existing)
    from lxml import etree

    effect_elem = etree.fromstring(soft_edge_xml)
    sp_pr.append(effect_elem)


def get_cyrillic_font():
    font_paths = [
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/calibri.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
    ]
    for path in font_paths:
        if os.path.exists(path):
            return path
    return None


def generate_image_via_text2image(prompt, access_token):
    url = "https://gigachat.devices.sberbank.ru/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": "GigaChat",
        "messages": [
            {
                "role": "user",
                "content": (
                    f"Создай реалистичное изображение на тему: {prompt}. "
                    "Без текста, без надписей, без искажений и артефактов. "
                    "Реалистичное расположение объектов: люди смотрят на мониторы, "
                    "мониторы повёрнуты экранами к людям."
                ),
            }
        ],
        "function_call": "auto",
    }
    for attempt in range(3):
        try:
            resp = requests.post(
                url, headers=headers, json=payload, timeout=90, verify=False
            )
            if resp.status_code == 200:
                data = resp.json()
                content = (
                    data.get("choices", [{}])[0].get("message", {}).get("content", "")
                )
                match = re.search(r"""<img\s+src\s*=\s*["']([^"']+)["']""", content)
                if match:
                    file_id = match.group(1)
                    file_url = f"https://gigachat.devices.sberbank.ru/api/v1/files/{file_id}/content"
                    file_resp = requests.get(file_url, headers=headers, verify=False)
                    if file_resp.status_code == 200:
                        return io.BytesIO(file_resp.content)
                else:
                    print(f"   Не найден тег <img> в ответе: {content[:100]}...")
            else:
                print(f"   text2image ответ {resp.status_code}: {resp.text[:100]}")
        except Exception as e:
            print(f"   Ошибка text2image (попытка {attempt + 1}): {e}")
            if attempt < 2:
                time.sleep(2)
    return None


def create_error_placeholder_image(width=600, height=400):
    img = Image.new("RGB", (width, height), color=(25, 50, 100))
    draw = ImageDraw.Draw(img)
    text = "К сожалению, изображение не удалось загрузить"
    font_path = get_cyrillic_font()
    if font_path:
        font = ImageFont.truetype(font_path, 24)
    else:
        font = ImageFont.load_default()
    lines = []
    words = text.split()
    line = ""
    for word in words:
        test_line = line + " " + word if line else word
        if draw.textbbox((0, 0), test_line, font=font)[2] <= width - 20:
            line = test_line
        else:
            lines.append(line)
            line = word
    if line:
        lines.append(line)
    y = 80
    for line in lines:
        tw = draw.textbbox((0, 0), line, font=font)[2]
        draw.text(((width - tw) // 2, y), line, fill=(255, 255, 255), font=font)
        y += 30
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def get_image_for_slide(slide_title, slide_content, access_token):
    if not ADD_IMAGES:
        return None
    prompt = slide_title
    if slide_content and len(slide_content) > 0:
        first_point = slide_content[0][:80]
        prompt = f"{slide_title}. {first_point}"
    if access_token:
        img = generate_image_via_text2image(prompt, access_token)
        if img:
            return img
    return create_error_placeholder_image()


def ask_giga(messages, max_tokens=1500, temperature=0.8):
    for attempt in range(1, 4):
        try:
            resp = giga.chat(
                {
                    "messages": messages,
                    "max_tokens": max_tokens,
                    "temperature": temperature,
                }
            )
            if resp and resp.choices:
                return resp.choices[0].message.content
            else:
                print(f"⚠️ Попытка {attempt}: пустой ответ")
        except Exception as e:
            print(f"⚠️ Попытка {attempt}: {e}")
            time.sleep(2)
    return None


def generate_intro_text(prompt, slides_titles):
    titles_str = ", ".join(slides_titles)
    sys_msg = (
        "Ты — ассистент, который пишет короткое введение к презентации. "
        "Напиши 1-2 предложения (не более 20 слов). Опиши важность темы и приведи ОДИН яркий факт или статистику. "
        "Не используй маркеры. Только связный текст."
    )
    user_msg = f"Тема презентации: «{prompt}». Слайды: {titles_str}"
    for attempt in range(3):
        text = ask_giga(
            [
                {"role": "system", "content": sys_msg},
                {"role": "user", "content": user_msg},
            ],
            max_tokens=200,
            temperature=0.7,
        )
        if text and len(text.strip()) > 10:
            return text.strip()
        time.sleep(1)
    return f"По данным исследований, тема «{prompt}» затрагивает каждого второго жителя страны. Рассмотрим ключевые аспекты."


def generate_conclusion_text(prompt, slides_info):
    titles_str = ", ".join(s["title"] for s in slides_info)
    key_points = []
    for s in slides_info:
        for point in s.get("content", [])[:2]:
            key_points.append(f"• {point}")
    points_str = "\n".join(key_points)
    sys_msg = (
        "Ты — ассистент, который пишет заключение к презентации. "
        "На основе темы и ключевых пунктов напиши 2-4 предложения. "
        "Они должны кратко подводить итог, давать практическую рекомендацию и логично завершать презентацию. "
        "Не используй маркеры. Только связный текст."
    )
    user_msg = f"Тема презентации: «{prompt}».\nЗаголовки слайдов: {titles_str}\nКлючевые пункты:\n{points_str}"
    for attempt in range(3):
        text = ask_giga(
            [
                {"role": "system", "content": sys_msg},
                {"role": "user", "content": user_msg},
            ],
            max_tokens=300,
            temperature=0.7,
        )
        if text and len(text.strip()) > 20:
            return text.strip()
        time.sleep(1)
    return f"Таким образом, соблюдение рассмотренных рекомендаций позволит минимизировать риски и повысить безопасность. Берегите себя!"


def generate_slide_structure(prompt, num_slides=5):
    dinstr = (
        "3–5 предложений с конкретными данными, цифрами, терминами. "
        "Пиши как эксперт в области безопасности, здоровья, полезных лайфхаков и методов борьбы с угрозами. "
        "Используй факты, статистику, практические рекомендации. "
        "Ответ должен быть максимально информативным и точным по теме."
    )

    print("   Шаг 1/2: заголовки...")
    sys1 = (
        f"Создай структуру: титульный слайд + ровно {num_slides} слайдов. "
        "Тематика: безопасность, здоровье, полезные лайфхаки, борьба с угрозами. "
        "JSON: {'title_slide': {'title':..., 'subtitle':...}, 'slides': [{'title':...}]}. "
        "Подзаголовок должен быть информативным (раскрывать суть: 'как избежать…', 'правила…', 'почему важно…'), "
        "а не просто повторять название. "
        "Никаких пояснений, только JSON."
    )
    ans1 = ask_giga(
        [
            {"role": "system", "content": sys1},
            {"role": "user", "content": f"Тема: {prompt}"},
        ],
        max_tokens=500,
    )
    if not ans1:
        return None

    ans1 = re.sub(r"```(?:json)?\s*(.*?)\s*```", r"\1", ans1, flags=re.DOTALL)
    ans1 = clean_json_string(ans1)
    data1 = parse_json_safe(ans1)
    if not data1 or "slides" not in data1:
        print("   Не удалось извлечь заголовки. Полный ответ (очищенный):", ans1[:300])
        return None

    slides_headers = data1["slides"]
    title_slide = data1.get("title_slide", {"title": prompt, "subtitle": ""})

    raw_sub = title_slide.get("subtitle", "")
    raw_sub = clean_subtitle(raw_sub)
    if raw_sub:
        raw_sub = normalize_title(raw_sub)
    else:
        raw_sub = ""

    if not raw_sub:
        sub = ask_giga(
            [
                {
                    "role": "user",
                    "content": (
                        f"Придумай ОДНУ короткую фразу (до 10 слов) – подзаголовок для презентации на тему: «{prompt}». "
                        "Подзаголовок должен раскрывать практическую пользу или главную идею: "
                        "например, «Как защитить себя и близких», «Правила безопасного поведения», "
                        "«Статистика и меры предосторожности». Только текст, без разметки."
                    ),
                }
            ],
            max_tokens=50,
        )
        if sub:
            raw_sub = sub.strip().strip('«»"')
            raw_sub = clean_subtitle(raw_sub)
            raw_sub = normalize_title(raw_sub)

    if not raw_sub:
        raw_sub = "Практические советы и анализ рисков"

    title_slide["subtitle"] = raw_sub

    print(f"   Шаг 2/2: контент для {len(slides_headers)} слайдов...")
    full_slides = []
    for idx, info in enumerate(slides_headers, 1):
        t = info.get("title", f"Слайд {idx}")
        print(f"      {idx}/{len(slides_headers)}: {t}")
        sys2 = (
            f"Тема: «{prompt}». Заголовок: «{t}». "
            f"Создай список из 3-5 экспертных пунктов, раскрывающих тему. "
            f"{dinstr} "
            "Верни СТРОГО JSON: {'content': ['строка1', 'строка2', ...]}. "
            "Никаких пояснений, только JSON."
        )

        content = None
        for attempt in range(3):
            ans2 = ask_giga(
                [
                    {"role": "system", "content": sys2},
                    {"role": "user", "content": f"Контент для слайда: {t}"},
                ],
                max_tokens=800,
                temperature=0.7,
            )
            if ans2:
                print(f"         Ответ (попытка {attempt + 1}): {ans2[:100]}...")
                ans2 = clean_json_string(ans2)
                data2 = parse_json_safe(ans2)
                if (
                    data2
                    and isinstance(data2.get("content"), list)
                    and len(data2["content"]) > 0
                ):
                    content = data2["content"]
                    break
            else:
                print(f"         Попытка {attempt + 1}: пустой ответ от API")
            time.sleep(1)

        if content:
            full_slides.append({"title": t, "content": content})
        else:
            full_slides.append(
                {
                    "title": t,
                    "content": [
                        f"Ключевые аспекты темы «{t}».",
                        "Информация будет дополнена в финальной версии презентации.",
                        "Рекомендации уточняются.",
                    ],
                }
            )

    structure = {"title_slide": title_slide, "slides": full_slides}
    structure["title_slide"]["title"] = normalize_title(
        structure["title_slide"]["title"]
    )
    for s in structure["slides"]:
        s["title"] = normalize_title(s["title"])

    slide_titles = [s["title"] for s in structure["slides"]]
    intro_text = generate_intro_text(prompt, slide_titles)
    structure["intro_text"] = intro_text

    conclusion_text = generate_conclusion_text(prompt, structure["slides"])
    structure["conclusion_text"] = conclusion_text

    return structure


def create_pptx_from_template(
    structure, output_path, presenter_name="", presentation_date=None, access_token=None
):
    if not os.path.exists(TEMPLATE_PATH):
        print("Ошибка: шаблон не найден")
        return
    if presentation_date is None:
        presentation_date = date.today().strftime("%d.%m.%Y")

    prs = Presentation(TEMPLATE_PATH)
    delete_all_slides(prs)

    # Титульный слайд
    title_layout = prs.slide_layouts[TITLE_LAYOUT_INDEX]
    slide = prs.slides.add_slide(title_layout)
    ts = structure.get("title_slide", {})
    ttext, stext = ts.get("title", ""), ts.get("subtitle", "")
    if len(ttext) > 70:
        tf, sf = Pt(24), Pt(14)
    else:
        tf, sf = None, None
    if slide.shapes.title:
        slide.shapes.title.text = ttext
        if tf:
            for p in slide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = tf
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = stext
            if sf:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = sf
            break
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 10:
            shape.text = presenter_name if presenter_name else ""
        elif shape.placeholder_format.idx == 11:
            shape.text = presentation_date

    # Вводный слайд
    intro_text = structure.get("intro_text", "")
    content_layout = prs.slide_layouts[CONTENT_LAYOUT_INDEX]
    intro_slide = prs.slides.add_slide(content_layout)
    if intro_slide.shapes.title:
        intro_slide.shapes.title.text = "Введение"

    shapes_to_delete = []
    for shape in intro_slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 12:
            shapes_to_delete.append(shape)
    for shape in shapes_to_delete:
        intro_slide.shapes._spTree.remove(shape._element)

    left = Inches(1.0)
    top = Inches(1.6)
    width = prs.slide_width - Inches(2.0)
    height = Inches(1.2)
    textbox = intro_slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.add_paragraph()
    p.text = intro_text
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    pPr = p._p.get_or_add_pPr()
    from lxml import etree

    for child in list(pPr):
        if (
            child.tag.endswith("}buChar")
            or child.tag.endswith("}buAutoNum")
            or child.tag.endswith("}buFont")
            or child.tag.endswith("}buClr")
            or child.tag.endswith("}buSz")
            or child.tag.endswith("}buNone")
        ):
            pPr.remove(child)
    etree.SubElement(pPr, qn("a:buNone"))
    pPr.set("marL", "0")
    pPr.set("indent", "0")

    qr_size = Cm(5.5)
    qr_left = (prs.slide_width - qr_size) // 2
    qr_top = Inches(3.6)
    shape = intro_slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        qr_left,
        qr_top,
        qr_size,
        qr_size,
    )
    shape.fill.background()
    shape.line.color.rgb = RGBColor(0xFF, 0x8C, 0x00)
    shape.line.width = Pt(2.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "Место для вашего QR-кода с опросом"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.alignment = PP_ALIGN.CENTER
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    # Основные слайды
    half_width = prs.slide_width // 2
    margin = Inches(0.4)

    for sd in structure.get("slides", []):
        slide = prs.slides.add_slide(content_layout)
        slide_title = sd.get("title", "")
        slide_content = sd.get("content", [])
        if slide.shapes.title:
            slide.shapes.title.text = slide_title

        body_shape = None
        for shape in slide.placeholders:
            if (
                shape.placeholder_format.idx == 12
                and shape.placeholder_format.type == 2
            ):
                body_shape = shape
                break
        if body_shape is None:
            for shape in slide.placeholders:
                if (
                    shape.placeholder_format.type == 2
                    and shape.placeholder_format.idx != 0
                ):
                    body_shape = shape
                    break

        if body_shape:
            body_shape.left = margin
            body_shape.top = Inches(1.7)
            body_shape.width = half_width - margin * 2
            tf = body_shape.text_frame
            tf.clear()
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for point in slide_content:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.font.size = Pt(14)
                p.space_after = Pt(6)

        img_stream = get_image_for_slide(slide_title, slide_content, access_token)
        if img_stream:
            try:
                pil_img = Image.open(img_stream)
                aspect = pil_img.width / pil_img.height

                top_margin = Inches(1.2)
                bottom_margin = Inches(0.8)
                max_img_height = Emu(prs.slide_height - top_margin - bottom_margin)
                max_img_width = half_width - margin * 2

                if max_img_width / aspect <= max_img_height:
                    img_width = max_img_width
                    img_height = Emu(int(img_width / aspect))
                else:
                    img_height = max_img_height
                    img_width = Emu(int(img_height * aspect))

                left = Emu(half_width + margin)
                top = Emu(top_margin)
                picture = slide.shapes.add_picture(
                    img_stream, left, top, width=img_width, height=img_height
                )

                add_soft_edge(picture, SOFT_EDGE_RADIUS)
                print("   ✓ изображение добавлено (с мягкими краями)")
            except Exception as e:
                print(f"   × ошибка вставки: {e}")

    # Заключительный слайд
    conclusion_text = structure.get("conclusion_text", "")
    conclusion_slide = prs.slides.add_slide(content_layout)
    if conclusion_slide.shapes.title:
        conclusion_slide.shapes.title.text = "Заключение"

    shapes_to_delete = []
    for shape in conclusion_slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 12:
            shapes_to_delete.append(shape)
    for shape in shapes_to_delete:
        conclusion_slide.shapes._spTree.remove(shape._element)

    left = Inches(1.0)
    top = Inches(1.8)
    width = prs.slide_width - Inches(2.0)
    height = Inches(2.5)
    textbox = conclusion_slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.add_paragraph()
    p.text = conclusion_text
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    pPr = p._p.get_or_add_pPr()
    for child in list(pPr):
        if (
            child.tag.endswith("}buChar")
            or child.tag.endswith("}buAutoNum")
            or child.tag.endswith("}buFont")
            or child.tag.endswith("}buClr")
            or child.tag.endswith("}buSz")
            or child.tag.endswith("}buNone")
        ):
            pPr.remove(child)
    etree.SubElement(pPr, qn("a:buNone"))
    pPr.set("marL", "0")
    pPr.set("indent", "0")

    static_line1 = (
        "Данная презентация сгенерирована искусственным интеллектом, "
        "текст и изображения могут иметь очевидные ошибки, но это отличный повод "
        "разбавить обстановку серьёзной встречи, вспомнить, что ИИ нас, людей, не заменит, "
        "и тепло улыбнуться коллегам."
    )
    static_line2_prefix = (
        "По вопросам улучшения генерации и исправления ошибок обращайтесь к "
    )
    static_email = "Aleksander.Klimov@evraz.com"

    static_left = Inches(1.0)
    static_top = Inches(6.0)
    static_width = prs.slide_width - Inches(2.0)
    static_height = Inches(1.0)
    textbox_static = conclusion_slide.shapes.add_textbox(
        static_left, static_top, static_width, static_height
    )
    tf_static = textbox_static.text_frame
    tf_static.word_wrap = True
    tf_static.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p1 = tf_static.add_paragraph()
    p1.text = static_line1
    p1.font.size = Pt(10)
    p1.alignment = PP_ALIGN.CENTER
    pPr1 = p1._p.get_or_add_pPr()
    for child in list(pPr1):
        if (
            child.tag.endswith("}buChar")
            or child.tag.endswith("}buAutoNum")
            or child.tag.endswith("}buFont")
            or child.tag.endswith("}buClr")
            or child.tag.endswith("}buSz")
            or child.tag.endswith("}buNone")
        ):
            pPr1.remove(child)
    etree.SubElement(pPr1, qn("a:buNone"))

    p2 = tf_static.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run_prefix = p2.add_run()
    run_prefix.text = static_line2_prefix
    run_prefix.font.size = Pt(10)
    run_email = p2.add_run()
    run_email.text = static_email
    run_email.font.size = Pt(10)
    run_email.hyperlink.address = f"mailto:{static_email}"
    run_suffix = p2.add_run()
    run_suffix.text = " "
    run_suffix.font.size = Pt(10)

    pPr2 = p2._p.get_or_add_pPr()
    for child in list(pPr2):
        if (
            child.tag.endswith("}buChar")
            or child.tag.endswith("}buAutoNum")
            or child.tag.endswith("}buFont")
            or child.tag.endswith("}buClr")
            or child.tag.endswith("}buSz")
            or child.tag.endswith("}buNone")
        ):
            pPr2.remove(child)
    etree.SubElement(pPr2, qn("a:buNone"))

    # Нумерация слайдов
    slide_number = 2
    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            continue
        width = Inches(0.8)
        height = Inches(0.4)
        left = prs.slide_width - width - Inches(0.4)
        top = prs.slide_height - height - Inches(0.3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.add_paragraph()
        p.text = str(slide_number)
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.RIGHT
        slide_number += 1

    prs.save(output_path)
    print(f"✅ Презентация сохранена: {output_path}")


# ─── Flask-приложение ─────────────────────────────────────────────
app = Flask(__name__)

HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="ru">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Генератор презентаций</title>
    <style>
        * { box-sizing: border-box; margin: 0; padding: 0; }
        body {
            font-family: 'Segoe UI', Arial, sans-serif;
            background: linear-gradient(135deg, #1a3c6e 0%, #2a5f8f 100%);
            min-height: 100vh;
            display: flex;
            justify-content: center;
            align-items: center;
            padding: 20px;
        }
        .container {
            background: white;
            border-radius: 16px;
            box-shadow: 0 20px 60px rgba(0,0,0,0.3);
            max-width: 500px;
            width: 100%;
            padding: 40px 30px;
        }
        h1 {
            color: #1a3c6e;
            font-size: 24px;
            margin-bottom: 10px;
            text-align: center;
        }
        p.subtitle {
            color: #666;
            font-size: 14px;
            margin-bottom: 30px;
            text-align: center;
        }
        label {
            display: block;
            font-weight: 600;
            color: #333;
            margin-bottom: 6px;
            font-size: 13px;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }
        input {
            width: 100%;
            padding: 12px 16px;
            margin-bottom: 20px;
            border: 2px solid #e0e0e0;
            border-radius: 8px;
            font-size: 15px;
            transition: border-color 0.3s;
            font-family: inherit;
        }
        input:focus {
            outline: none;
            border-color: #1a3c6e;
        }
        button {
            width: 100%;
            padding: 14px;
            background: #1a3c6e;
            color: white;
            border: none;
            border-radius: 8px;
            font-size: 16px;
            font-weight: 600;
            cursor: pointer;
            transition: background 0.3s, transform 0.2s;
        }
        button:hover {
            background: #0f2b4f;
            transform: translateY(-1px);
        }
        button:active {
            transform: translateY(0);
        }
        #status {
            display: none;
            text-align: center;
            margin-top: 30px;
            padding: 20px;
            background: #f8f9fa;
            border-radius: 12px;
        }
        .spinner {
            border: 4px solid #e0e0e0;
            border-top: 4px solid #1a3c6e;
            border-radius: 50%;
            width: 48px;
            height: 48px;
            animation: spin 0.8s linear infinite;
            margin: 0 auto 16px;
        }
        @keyframes spin {
            0% { transform: rotate(0deg); }
            100% { transform: rotate(360deg); }
        }
        #message {
            color: #555;
            font-size: 15px;
            margin-bottom: 12px;
        }
        #downloadLink {
            display: none;
            color: #1a3c6e;
            font-weight: 600;
            text-decoration: none;
            font-size: 16px;
            padding: 10px 20px;
            background: #e8f0fe;
            border-radius: 8px;
            transition: background 0.3s;
        }
        #downloadLink:hover {
            background: #d0e2ff;
        }
    </style>
</head>
<body>
    <div class="container">
        <h1>📊 Генератор презентаций</h1>
        <p class="subtitle">Создайте профессиональную презентацию с помощью GigaChat AI</p>
        <form id="generateForm">
            <label for="theme">Тема презентации</label>
            <input type="text" id="theme" placeholder="Например: Безопасность на воде в летний период" required>

            <label for="slides">Количество слайдов</label>
            <input type="number" id="slides" value="5" min="1" max="15">

            <label for="presenter">ФИО докладчика</label>
            <input type="text" id="presenter" placeholder="Иванов Иван Иванович">

            <button type="submit">✨ Сгенерировать презентацию</button>
        </form>

        <div id="status">
            <div class="spinner" id="spinner"></div>
            <p id="message">Идёт генерация презентации...</p>
            <a id="downloadLink" href="#">📥 Скачать презентацию</a>
        </div>
    </div>

    <script>
        document.getElementById('generateForm').onsubmit = async function(e) {
            e.preventDefault();
            const status = document.getElementById('status');
            const spinner = document.getElementById('spinner');
            const message = document.getElementById('message');
            const downloadLink = document.getElementById('downloadLink');

            status.style.display = 'block';
            spinner.style.display = 'block';
            message.textContent = 'Генерация структуры и изображений...';
            downloadLink.style.display = 'none';

            const formData = new FormData();
            formData.append('theme', document.getElementById('theme').value.trim());
            formData.append('slides', document.getElementById('slides').value);
            formData.append('presenter', document.getElementById('presenter').value.trim());

            try {
                const response = await fetch('/generate', {
                    method: 'POST',
                    body: formData
                });

                if (!response.ok) {
                    const err = await response.json();
                    throw new Error(err.error || 'Неизвестная ошибка сервера');
                }

                const blob = await response.blob();
                const url = window.URL.createObjectURL(blob);
                downloadLink.href = url;
                downloadLink.download = 'presentation.pptx';
                downloadLink.style.display = 'inline-block';
                message.textContent = '✅ Презентация готова!';
            } catch (error) {
                message.textContent = '❌ Ошибка: ' + error.message;
            } finally {
                spinner.style.display = 'none';
            }
        };
    </script>
</body>
</html>
"""


@app.route("/")
def index():
    return render_template_string(HTML_TEMPLATE)


@app.route("/generate", methods=["POST"])
def generate():
    theme = request.form.get("theme", "").strip()
    if not theme:
        return jsonify({"error": "Тема не может быть пустой"}), 400

    try:
        num_slides = int(request.form.get("slides", "5"))
        num_slides = max(1, min(15, num_slides))
    except ValueError:
        num_slides = 5

    presenter = request.form.get("presenter", "").strip()

    try:
        # Генерация структуры
        structure = generate_slide_structure(theme, num_slides)
        if not structure:
            return jsonify({"error": "Не удалось создать структуру презентации"}), 500

        # Создаём временный файл
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            temp_path = tmp.name

        # Генерируем презентацию
        access_token = giga.token if hasattr(giga, "token") else None
        create_pptx_from_template(
            structure, temp_path, presenter_name=presenter, access_token=access_token
        )

        # Читаем файл в память и отправляем
        with open(temp_path, "rb") as f:
            pptx_data = f.read()
        os.unlink(temp_path)  # удаляем временный файл

        pptx_buffer = io.BytesIO(pptx_data)
        pptx_buffer.seek(0)

        filename = safe_filename(theme) + ".pptx"
        return send_file(
            pptx_buffer,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=filename,
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    # Локальный запуск: python app.py
    import sys

    if len(sys.argv) > 1 and sys.argv[1] == "--cloud":
        # Запуск через waitress на Render
        port = int(os.environ.get("PORT", 8080))
        print(f"Запуск production сервера на порту {port}")
        waitress.serve(app, host="0.0.0.0", port=port)
    else:
        print("=" * 60)
        print("🚀 Сервер генератора презентаций запущен")
        print(f"📂 Шаблон презентации: {TEMPLATE_PATH}")
        print("🌐 Откройте в браузере: http://localhost:5000")
        print("=" * 60)
        app.run(debug=False, host="0.0.0.0", port=5000)
